# core.py - public API
# coding=utf-8

from pptx import Presentation

import sys
import codecs
import logging

import pptx_template.pptx_util as util
import pptx_template.text as txt
import pptx_template.chart as ch
import pptx_template.table as tb

log = logging.getLogger()

def edit_slide(slide, model, skip_model_not_found = False):
    """
        1つのスライドに対して文字列置換およびチャートCSV設定を行う
        チャート設定や文字列置換は1スライドに対して複数持てる。配列やdictなどで引渡し、pptxからはEL式で特定する
        チャート設定なのか文字列置換なのかは、EL式の配置されているpptx内のオブジェクトで判断される
        文字列置換やチャートタイトルに EL式で {answer.0} のような形で指定する
        チャート設定として指定可能な項目:
         - file_name : CSVファイルの名前
         - body: CSVファイルの中身そのものを直接文字列で指定できる(file_name, file_encodingは無視される)
         - value_axis_max: チャート左側軸の最大値。(省略可)
         - value_axix_min: チャート左側軸の最小値。(省略可)
         - (TBI) file_encoding: CSVファイルのエンコーディング。省略時は utf-8
    """

    # pptx内の TextFrame の EL表記を model の値で置換する
    for shape in txt.select_all_text_shapes(slide):
        try:
            txt.replace_all_els_in_text_frame(shape.text_frame, model)
        except:
            if not skip_model_not_found:
                raise
    for shape in txt.select_all_tables(slide):
        try:
            txt.replace_all_els_in_table(shape, model)
        except:
            if not skip_model_not_found:
                raise

    # pptx内の 各チャートに対してcsvの値を設定する
    for chart in ch.select_all_chart_shapes(slide):
        try:
            ch.load_data_into_chart(chart, model)
        except:
            if not skip_model_not_found:
                raise

    # pptx内の 各テーブルに対してcsvの値を設定する
    for shape in txt.select_all_tables(slide):
        try:
            tb.load_data_into_table(shape, model)
        except:
            if not skip_model_not_found:
                raise

def remove_slide(presentation, slide):
    """
     presentation から 指定した slide を削除する
    """
    util.remove_slide(presentation, slide)


def remove_slide_id(presentation, slide_id):
    """
         指定した id のスライドから {id:foobar} という形式の文字列を削除する
    """
    slide = get_slide(presentation, slide_id)
    for shape in txt.select_all_text_shapes(slide):
        if txt.extract_slide_id(shape.text) == slide_id:
            shape.text = ''

def remove_all_slides_having_id(presentation):
    """
         {id:foobar} という文字列を持つすべてのスライドを削除する
    """
    unused_slides = []
    for slide in presentation.slides:
        for shape in txt.select_all_text_shapes(slide):
            slide_id = txt.extract_slide_id(shape.text)
            if slide_id:
                unused_slides.append((slide_id, slide))
                break
    for slide_id, slide in unused_slides:
        log.info("Removing unused slide_id: %s" % slide_id)
        remove_slide(presentation, slide)


def get_slide(presentation, slide_id):
    """
         指定した id に対して {id:foobar} という TextFrame を持つスライドを探す
    """
    for slide in presentation.slides:
        for shape in txt.select_all_text_shapes(slide):
            if txt.extract_slide_id(shape.text) == slide_id:
                return slide
    raise ValueError(u"slide id:%s not found" % slide_id)
