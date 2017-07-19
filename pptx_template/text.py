#
# coding=utf-8

import logging
import re
import numbers

from six import string_types

import pptx_template.pyel as pyel

from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.table import Table

log = logging.getLogger()

_EL_RE = re.compile(r"\{([A-Za-z0-9._\-]+)\}")
_SLIDE_ID_RE = re.compile(r"\{id:([A-Za-z0-9._\-]+)\}")

def extract_slide_id(text):
    match = _SLIDE_ID_RE.search(text)
    if match:
        return match.group(1)
    else:
        return None

def _iterate_els(text):
    pos = 0
    while pos < len(text):
        text_id_match = _EL_RE.search(text[pos:])
        if text_id_match:
            pos = pos + text_id_match.end(1) + 1
            yield text_id_match.group(1)
        else:
            break;

def _el_to_placeholder(el):
    return u"{%s}" % el


def search_first_el(text):
    for id in _iterate_els(text):
        return id
    return None


def select_all_text_shapes(slide):
    return [ s for s in slide.shapes if s.shape_type in [1,14,17] ]


def select_all_tables(slide):
    return [ s.table for s in slide.shapes if isinstance(s, GraphicFrame) and s.shape_type == 19 ]


def replace_all_els_in_table(table, model):
    """
     table の各セルの中に EL 形式があれば、それを model の該当する値と置き換える
    """
    for cell in [ cell for row in table.rows for cell in row.cells ]:
        replace_all_els_in_text_frame(cell.text_frame, model)


def replace_el_in_text_frame_with_str(text_frame, el, replacing_text):
    """
     text_frame の各 paragraph.run 中のテキストに指定の EL 形式があれば、それを replacing_text で置き換える
    """
    placeholder = _el_to_placeholder(el)
    for paragraph in text_frame.paragraphs:
        if not placeholder in paragraph.text:
            continue

        original_run_for_debug_log = [r.text for r in paragraph.runs]
        ((start_run, start_pos), (end_run, end_pos)) = _find_el_position([r.text for r in paragraph.runs], el)
        for (i, run) in enumerate(paragraph.runs):
            if i == start_run and i == end_run:
                run.text = run.text.replace(placeholder, replacing_text)
            elif i == start_run:
                run.text = run.text[0:start_pos] + replacing_text
            elif i == end_run:
                run.text = run.text[end_pos + 1:]
                break
            elif start_run < i and i < end_run:
                run.text = ''
        log.debug(u" Replaced text: %s  --> %s" % (original_run_for_debug_log, [r.text for r in paragraph.runs]))
        return True
    return False

def replace_all_els_in_text_frame(text_frame, model):
    """
     text_frame 中のテキストに EL 形式が一つ以上あれば、それを model の該当する値と置き換える
    """
    for el in _iterate_els(text_frame.text):
        value = pyel.eval_el(el, model)
        if not value:
            replacing_text = ''
        elif isinstance(value, numbers.Number):
            replacing_text = str(value)
        elif not isinstance(value, string_types):
            log.error(u"Invalid value for {%s}, model: %s" % (el, value))
            continue
        else:
            replacing_text = value

        if not replace_el_in_text_frame_with_str(text_frame, el, replacing_text):
            log.error(u"Cannot find {%s} in one text-run. To fix this, select this whole EL [%s] and reset font size by clicking size up then down" % (text_id, text_frame.text))

def _find_el_position(texts, el):
    """
    text の配列中に分かれて記述されている EL の、先頭の '\{' の位置と、最後の '\}' の位置を返す。
    それぞれの位置は (text_index, position_in_text) の形で返される。
    """
    placeholder = _el_to_placeholder(el)
    full_text = ''.join(texts)

    start_pos = full_text.find(placeholder)
    if start_pos < 0:
        raise ValueError(u"texts %s doesn't contain EL:%s" % (texts, el))
    end_pos = start_pos + len(placeholder) - 1

    start_run_pos = start_run_index = end_run_pos = end_run_index = -1
    for (run_index, text) in enumerate(texts):
        length = len(text)

        # log.error("run_index:%(run_index)s start_pos:%(start_pos)s end_pos:%(end_pos)s length:%(length)s" % locals())
        if start_pos >= 0 and start_pos < length:
            start_run_index = run_index
            start_run_pos = start_pos
        start_pos = start_pos - length

        if end_pos >= 0 and end_pos < length:
            end_run_index = run_index
            end_run_pos = end_pos
            break
        end_pos = end_pos - length

    return ((start_run_index, start_run_pos), (end_run_index, end_run_pos))
