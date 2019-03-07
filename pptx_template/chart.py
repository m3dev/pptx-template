#
# coding=utf-8

import logging
import os.path
import math
from io import StringIO

from pptx.shapes.graphfrm import GraphicFrame
from pptx.chart.data import ChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE as ct
from pptx.chart.chart import Chart

import pandas as pd
import numpy as np

from six import string_types

import pptx_template.pyel as pyel
import pptx_template.text as txt
import pptx_template.pptx_util as util

log = logging.getLogger()

def _nan_to_none(x):
    # log.debug(u" type of x:%s is:%s" % (x, type(x)))
    if isinstance(x, np.generic):
        result = None if np.isnan(x) else x.item()
    elif isinstance(x, string_types):
        result = _to_unicode(x)
    elif math.isnan(x):
        result = None
    else:
        result = x
    return result

def _to_unicode(s):
    return s if isinstance(s, type(u"a")) else unicode(s,'utf-8')

def _build_xy_chart_data(csv, xy_transpose, number_format):
    chart_data = XyChartData()
    for i in range(1, csv.columns.size):
        # nameに日本語が入ると後続処理中で、python v2.7の場合にUnicodeDecodeErrorが出るため対処。nameは結局pptx内では使われない
        series = chart_data.add_series(u"column%s" % i, number_format=_normalize_number_format(number_format))
        xy_col = csv.iloc[:, [0, i]]
        for (_, row) in xy_col.iterrows():
            if xy_transpose:
                x, y = _nan_to_none(row[0]), _nan_to_none(row[1])
            else:
                y, x = _nan_to_none(row[0]), _nan_to_none(row[1])
            log.debug(u" Adding xy %s,%s" % (x, y))
            series.add_data_point(x, y)
    return chart_data

def _build_chart_data(csv, number_format):
    chart_data = ChartData()
    categories = [_nan_to_none(x) or '' for x in csv.iloc[:,0].values]
    log.debug(u" Setting categories with values:%s" % categories)
    chart_data.categories = categories

    for i in range(1, csv.columns.size):
        col = csv.iloc[:, i]
        values = [_nan_to_none(x) for x in col.values]
        name = _to_unicode(col.name)
        log.debug(u" Adding series:%s values:%s" % (name, values))
        # 本来、number_formatは既存のchartの設定をそのまま引き継ぎたかったが、
        # python-pptx v0.6.17 では、既存のchartのchart_dataを取得するAPIは存在せず、
        # 新たにchart_dataを作って、chart.replace_data() する必要がある。
        # そのため、number_formatは、modelのoptionから取得する方針とする。
        chart_data.add_series(name, values, _normalize_number_format(number_format))
    return chart_data

def _normalize_number_format(number_format):
    """
    pptx.chartの内部で、'\\'が’\'扱いになってしまうため、更にエスケープ処理を行う
    """
    return number_format.replace('\\','\\\\') if number_format != None else number_format

def _is_xy_chart(chart):
    xy_charts = [ct.XY_SCATTER_LINES, ct.XY_SCATTER_LINES_NO_MARKERS, ct.XY_SCATTER, ct.XY_SCATTER_SMOOTH, ct.XY_SCATTER_SMOOTH_NO_MARKERS]
    return chart.chart_type in xy_charts

def _set_value_axis(chart, chart_id, chart_setting):
    max = chart_setting.get('value_axis_max')
    min = chart_setting.get('value_axis_min')
    if max or min:
        util.set_value_axis(chart, max = max, min = min)

def _load_csv_into_dataframe(chart_id, chart_setting):
    if 'body' in chart_setting:
        csv_body = chart_setting.get('body')
        return pd.read_csv(StringIO(csv_body), index_col=False)
    elif 'tsv_body' in chart_setting:
        tsv_body = chart_setting.get('tsv_body')
        return pd.read_csv(StringIO(tsv_body), delimiter='\t', index_col=False)
    else:
        csv_file_name = chart_setting.get('file_name')
        if not csv_file_name:
            for ext in ['csv', 'tsv']:
                csv_file_name = "%s.%s" % (chart_id, ext)
                if os.path.isfile(csv_file_name):
                    break
            else:
                raise ValueError(u"File not found: csv or tsv for %s" % chart_id)

        log.debug(u" Loading from csv file: %s" % csv_file_name)
        delimiter = '\t' if csv_file_name.endswith('.tsv') else ','
        return pd.read_csv(csv_file_name, delimiter=delimiter, index_col=False)

def _replace_chart_data_with_csv(chart, chart_id, chart_setting):
    """
        1つのチャートに対して指定されたCSVからデータを読み込む。
    """
    log.debug(chart_setting)
    csv = _load_csv_into_dataframe(chart_id, chart_setting)
    log.debug(u" Loaded Data:\n%s" % csv)

    number_format = chart_setting.get("number_format")
    xy_transpose = chart_setting.get("xy_transpose")
    if _is_xy_chart(chart):
        log.info(u"Setting csv/tsv into XY chart_id: %s" % chart_id)
        chart_data = _build_xy_chart_data(csv, xy_transpose, number_format)
    else:
        log.info(u"Setting csv/tsv into chart_id: %s" % chart_id)
        chart_data = _build_chart_data(csv, number_format)

    chart.replace_data(chart_data)

    log.debug(u" Completed chart data replacement.")

    return


def load_data_into_chart(chart, model):
    log.debug(u"model:%s" % model)
    # チャートタイトルから {} で囲まれた文字列を探し、それをキーとしてチャート設定と紐付ける
    if not chart.has_title or not chart.chart_title.has_text_frame:
        return

    # チャートIDを取得
    title_frame = chart.chart_title.text_frame
    chart_id = txt.search_first_el(title_frame.text)
    if not chart_id:
        return

    # チャートタイトル中のチャートIDを削除
    txt.replace_el_in_text_frame_with_str(title_frame, chart_id, '')

    # チャートタイトル中のEL式の置換を行う
    txt.replace_all_els_in_text_frame(title_frame, model)

    chart_setting = pyel.eval_el(chart_id, model)
    log.debug(u" Found chart_id: %s, chart_setting: %s" % (chart_id, chart_setting))

    # チャートにデータを流し込む
    _replace_chart_data_with_csv(chart, chart_id, chart_setting)

    # 軸の最大値、最小値を設定
    _set_value_axis(chart, chart_id, chart_setting)

def select_all_chart_shapes(slide):
    return [ s.chart for s in slide.shapes if isinstance(s, GraphicFrame) and s.shape_type == 3 ]
